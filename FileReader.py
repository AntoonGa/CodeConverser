# -*- coding: utf-8 -*-
"""
Created on Tue Aug  8 15:47:38 2023

@author: agarc

"""
from PathExplorer import PathExplorer
from PyPDF2 import PdfReader
import docx
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx import Presentation
import pandas as pd
from tqdm import tqdm
# =============================================================================
# FileReader (inherits from PathExplorer)
# =============================================================================


class FileReader(PathExplorer):
    """
    This class takes care of extracting text from a single document.
    /!\ This is where most of the work must go:
        Add or improve documents reader functions here.
        Ideally we would make sure to keep "metadata" when reading files, such as "Title, Section, Tab" etC.
    """

    def __init__(self):
        super().__init__()

        # list of valid extentions for which a loader is ready.
        # should add: .doc, .ppt
        self.loader_router = {'.pptx': self._read_text_from_pptx,
                              '.docx': self._read_text_from_docx,
                              '.pdf': self._read_text_from_pdf,
                              '.txt': self._read_text_from_txt,
                              '.py': self._read_text_from_txt,
                              '.md': self._read_text_from_txt}

        return

    def get_all_documents_in_df(self, directory: str, read_only_extensions=[], ignore_extensions=[]):
        """
        Reads all documents in the specified directory and returns a DataFrame containing
        the document information and extracted text for each document.

        Args:
            directory (str): The path to the directory containing the documents.

        Returns:
            pd.DataFrame: A DataFrame containing the document information and extracted text for each document.
        """
        # Get all file paths in the directory
        file_paths = self.get_all_files_paths(directory)

        # filter out unwanted or wanted extension:
        file_paths = self._filter_extensions(
            file_paths, read_only_extensions=read_only_extensions, ignore_extensions=ignore_extensions)

        # Initialize an empty list of hashmap to store the document information
        all_documents = []
        # Loop through the file paths and read each document
        for file_path in tqdm(file_paths):
            document_info = self._read_single_document(file_path)
            if document_info:
                all_documents.append(document_info)

        df_all_documents = pd.DataFrame(all_documents)
        return df_all_documents

    def _read_single_document(self, file_path) -> dict:
        """
        Reads a single documents.
        Uses the exention to send to the appropriate text exactor

        /!\ returns a hashmap /!\ 

        Parameters
        ----------
        file_path : file path

        Returns
        -------
        hashmap : {"file_path" : file_path,
                    "file_name" : file_name,
                    "file_extension" : file_extension,
                    "file_text" : text}
        """

        # grab extension and name from file name
        file_extension = self._get_single_file_extension(file_path)
        file_name = self._get_single_file_name(file_path)

        # check if the file exitsts or if exention is handled by the router
        if self._assert_file_exists(file_path) == False or file_extension not in self.loader_router:
            source_text = {"file_path": file_path,
                           "file_name": file_name,
                           "file_extension": file_extension,
                           "file_full_name": file_name + file_extension,
                           "file_text": "This file does not exist"}
            print('Not readable format ', file_extension)
        # if the extension is valid:
        else:
            # exctract the text
            text = self.loader_router[file_extension](file_path)
            # dump into a hashmap
            source_text = {"file_path": file_path,
                           "file_name": file_name,
                           "file_extension": file_extension,
                           "file_full_name": file_name + file_extension,
                           "file_text": text}

        return source_text

    def _read_text_from_pptx(self, file_path: str) -> str:
        """
        Read all text from a PowerPoint (PPTX) file.
        Args:
        file_path (str): The path of the PPTX file.
        Returns:
        str: The extracted text from the PPTX file.

        Dependencies:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE        
        """

        def extract_text_from_shape(shape):
            """
            Extract any text from mutiple shapes types
            This could be defined recursively also but this seems easyier that way.
            """
            text_runs = []
            if shape.has_text_frame:  # from text frame
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:  # from tables
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                text_runs.append(run.text)
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:  # from groups
                for sub_shape in shape.shapes:
                    text_runs.extend(extract_text_from_shape(sub_shape))
            elif shape.shape_type == MSO_SHAPE_TYPE.MEDIA:  # from media
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text_runs.append(run.text)

            return " ".join(text_runs).strip()

        # load presentation
        prs = Presentation(file_path)
        # instantiate empty output
        text_runs = []
        # loop through slides and slide.shapes
        for slide in prs.slides:
            for shape in slide.shapes:
                # get text from shape
                shape_text = extract_text_from_shape(shape)
                if shape_text:
                    text_runs.append(shape_text)

        # return only if not empty
        if text_runs:
            text = "\n".join(text_runs)
            return text
        else:
            return " "

    def _read_text_from_docx(self, file_path: str) -> str:
        """
        Read all text from a word (docs) file.
        Args:
        file_path (str): The path of the word file.
        Returns:
        str: The extracted text from the word file.

        Dependencies:
        import docx       
        """

        # load document
        doc = docx.Document(file_path)
        # instantiate return
        full_text = []

        # parse through paragraphs
        for paragraph in doc.paragraphs:
            paragraph_text = paragraph.text.strip()
            if paragraph_text:
                full_text.append(paragraph.text)

        # parse through tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        paragraph_text = paragraph.text.strip()
                        if paragraph_text:
                            full_text.append(paragraph_text)
            if full_text:
                full_text.append('_____\n')

        # return only if not empty
        if full_text:
            text = '\n'.join(full_text)
            return text
        else:
            return " "

    def _read_text_from_pdf(self, file_path: str) -> str:
        """
        Read all text from a pdf file.
        Args:
        file_path (str): The path of the pdf file.
        Returns:
        str: The extracted text from the pdf file.

        Dependencies:
        from PyPDF2 import PdfReader
        """

        # creating a pdf reader object
        reader = PdfReader(file_path)

        # instantiate return
        text = ''
        # getting a specific page from the pdf file
        for page in reader.pages:
            # extracting text from page
            page_text = page.extract_text().strip()
            if page_text:
                text = text + page_text

        # return only if not empty
        if text:
            text = text
            return text
        else:
            return " "

    def _read_text_from_txt(self, file_path: str) -> str:
        """
        Read all text from a .txt fie

        Parameters
        ----------
        file_path : str
            the path of the txt file.

        Returns
        -------
        str
            the extracted text from the txt file.

        """
        with open(file_path, 'r', encoding='utf-8') as file:
            text = file.read()

        if text:
            return text
        else:
            return " "


if __name__ == "__main__":
    reader = FileReader()
    directory = r'../_data_raw/'
    text_dataframe = reader.get_all_documents_in_df(directory, read_only_extensions=['.pptx'], ignore_extensions=[])
